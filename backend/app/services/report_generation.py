"""Board-style PDF and PowerPoint report generation -- the same Income
Statement/Balance Sheet/KPI data the Excel exports already use, formatted
as an actual document instead of a raw table dump: a title, section
headings, and styled tables, suitable to hand to someone who isn't going
to open a spreadsheet.
"""
from dataclasses import dataclass
from datetime import date
from io import BytesIO
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt
from reportlab.lib import colors
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import inch
from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle
from sqlalchemy.orm import Session

from app.models import Company
from app.services import financial_statements, kpis

BRAND_COLOR = "#2f5d50"


@dataclass
class BoardReportData:
    company_name: str
    currency: str
    start_period: date
    end_period: date
    as_of: date
    kpi: kpis.KPISet
    income_statement: financial_statements.IncomeStatement
    balance_sheet: financial_statements.BalanceSheet


def build_report_data(db: Session, company_id, start_period: date, end_period: date, as_of: date) -> BoardReportData:
    company = db.get(Company, company_id)
    kpi = kpis.compute_kpis(db, company_id, fiscal_year=start_period.year)
    income = financial_statements.income_statement(db, company_id, start_period, end_period)
    balance = financial_statements.balance_sheet(db, company_id, as_of)
    return BoardReportData(
        company_name=company.name if company else "Unknown company",
        currency=company.base_currency if company else "USD",
        start_period=start_period,
        end_period=end_period,
        as_of=as_of,
        kpi=kpi,
        income_statement=income,
        balance_sheet=balance,
    )


def _fmt_amount(value: float) -> str:
    return f"{value:,.2f}"


def _fmt_pct(value: Optional[float]) -> str:
    return "—" if value is None else f"{value:.1f}%"


def _kpi_rows(kpi: kpis.KPISet) -> list[list[str]]:
    return [
        ["Metric", "Value"],
        ["Gross Margin", _fmt_pct(kpi.gross_margin_pct)],
        ["Budget Utilization", _fmt_pct(kpi.budget_utilization_pct)],
        ["Forecast Accuracy (MAPE)", _fmt_pct(kpi.forecast_accuracy_mape)],
        ["Cash Runway", f"{kpi.cash_runway_months} mo" if kpi.cash_runway_months is not None else "12+ mo"],
    ]


def _income_statement_rows(statement: financial_statements.IncomeStatement) -> list[list[str]]:
    rows = [["Account", "Amount"], ["Revenue", ""]]
    for line in statement.revenue_lines:
        rows.append([f"  {line.code} {line.name}", _fmt_amount(line.amount)])
    rows.append(["Total Revenue", _fmt_amount(statement.total_revenue)])
    rows.append(["Expenses", ""])
    for line in statement.expense_lines:
        rows.append([f"  {line.code} {line.name}", _fmt_amount(line.amount)])
    rows.append(["Total Expense", _fmt_amount(statement.total_expense)])
    rows.append(["Net Profit", _fmt_amount(statement.net_profit)])
    return rows


def _balance_sheet_rows(sheet: financial_statements.BalanceSheet) -> list[list[str]]:
    rows = [["Account", "Amount"], ["Assets", ""]]
    for line in sheet.asset_lines:
        rows.append([f"  {line.code} {line.name}", _fmt_amount(line.amount)])
    rows.append(["Total Assets", _fmt_amount(sheet.total_assets)])
    rows.append(["Liabilities", ""])
    for line in sheet.liability_lines:
        rows.append([f"  {line.code} {line.name}", _fmt_amount(line.amount)])
    rows.append(["Total Liabilities", _fmt_amount(sheet.total_liabilities)])
    rows.append(["Equity", ""])
    for line in sheet.equity_lines:
        rows.append([f"  {line.code} {line.name}", _fmt_amount(line.amount)])
    rows.append(["Total Equity", _fmt_amount(sheet.total_equity)])
    status = "Balanced" if sheet.is_balanced else f"Off by {_fmt_amount(abs(sheet.difference))}"
    rows.append([status, ""])
    return rows


def _pdf_table(rows: list[list[str]]) -> Table:
    table = Table(rows, colWidths=[3.5 * inch, 2.5 * inch])
    table.setStyle(
        TableStyle(
            [
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor(BRAND_COLOR)),
                ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                ("ALIGN", (1, 0), (1, -1), "RIGHT"),
                ("FONTSIZE", (0, 0), (-1, -1), 9),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
                ("TOPPADDING", (0, 0), (-1, -1), 4),
                ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#d0d0d0")),
                ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#f5f6f4")]),
            ]
        )
    )
    return table


def render_pdf(data: BoardReportData) -> bytes:
    buffer = BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=letter, topMargin=0.75 * inch, bottomMargin=0.75 * inch)
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle("ReportTitle", parent=styles["Title"], textColor=colors.HexColor(BRAND_COLOR))
    heading_style = ParagraphStyle("SectionHeading", parent=styles["Heading2"], textColor=colors.HexColor(BRAND_COLOR), spaceBefore=18)

    elements = [
        Paragraph(f"{data.company_name} — Financial Report", title_style),
        Paragraph(
            f"Income Statement: {data.start_period.isoformat()} to {data.end_period.isoformat()} "
            f"&middot; Balance Sheet as of {data.as_of.isoformat()} &middot; {data.currency}",
            styles["Normal"],
        ),
        Spacer(1, 16),
        Paragraph("Key Performance Indicators", heading_style),
        _pdf_table(_kpi_rows(data.kpi)),
        Paragraph("Income Statement", heading_style),
        _pdf_table(_income_statement_rows(data.income_statement)),
        Paragraph("Balance Sheet", heading_style),
        _pdf_table(_balance_sheet_rows(data.balance_sheet)),
    ]
    doc.build(elements)
    return buffer.getvalue()


def _add_table_slide(prs: Presentation, title: str, rows: list[list[str]]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title

    n_rows, n_cols = len(rows), len(rows[0])
    left, top, width, height = Inches(0.6), Inches(1.5), Inches(prs.slide_width.inches - 1.2), Inches(0.4 * n_rows)
    shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, min(height, Inches(5.5)))
    table = shape.table
    table.columns[0].width = Inches((prs.slide_width.inches - 1.2) * 0.65)
    for col in range(1, n_cols):
        table.columns[col].width = Inches((prs.slide_width.inches - 1.2) * 0.35 / (n_cols - 1))

    for row_idx, row in enumerate(rows):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.text = value
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(12 if row_idx == 0 else 11)
                paragraph.font.bold = row_idx == 0


def render_pptx(data: BoardReportData) -> bytes:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = data.company_name
    title_slide.placeholders[1].text = (
        f"Financial Report · {data.start_period.isoformat()} to {data.end_period.isoformat()} · {data.currency}"
    )

    _add_table_slide(prs, "Key Performance Indicators", _kpi_rows(data.kpi))
    _add_table_slide(prs, "Income Statement", _income_statement_rows(data.income_statement))
    _add_table_slide(prs, "Balance Sheet", _balance_sheet_rows(data.balance_sheet))

    buffer = BytesIO()
    prs.save(buffer)
    return buffer.getvalue()
